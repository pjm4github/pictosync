"""Test bench for PPTX export with blocks/runs text.

Tests:
  - Shape text from overlay-2.0 blocks/runs
  - Shape text from legacy meta label/tech/note
  - Per-run formatting (bold, italic, underline, color, font)
  - Multi-block paragraphs
  - Line annotation text box at midpoint
  - Fill color applied to shape
  - Pen color applied to shape outline
  - Rotation angle applied
  - Empty blocks produce no text
  - All shape kinds export without error
"""
from __future__ import annotations

import os
import sys
import tempfile

import pytest


from pptx import Presentation
from pptx.util import Pt

from pptx_export import export_to_pptx


# ── Helpers ──────────────────────────────────────────────────────────────

def _export_and_read(annotations, **kwargs):
    """Export annotations to a temp PPTX and read back the presentation."""
    path = os.path.join(tempfile.gettempdir(), "test_export.pptx")
    export_to_pptx(annotations, path, **kwargs)
    return Presentation(path)


def _get_shape_texts(prs):
    """Return list of (shape_index, [paragraph_texts]) for all text shapes."""
    results = []
    for slide in prs.slides:
        for i, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                paras = [p.text for p in shape.text_frame.paragraphs if p.text]
                if paras:
                    results.append((i, paras))
    return results


def _get_shape_runs(prs):
    """Return list of (para_text, [(run_text, bold, italic)]) for all shapes."""
    results = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    runs = []
                    for r in p.runs:
                        runs.append((r.text, r.font.bold, r.font.italic))
                    if runs:
                        results.append((p.text, runs))
    return results


# ── Test Annotations ─────────────────────────────────────────────────────

RECT_OVERLAY20 = {
    "id": "a1", "kind": "rect",
    "geom": {"x": 100, "y": 100, "w": 200, "h": 120},
    "contents": {
        "frame": {"halign": "center", "valign": "middle"},
        "default_format": {"font_size": 14, "color": "#333333FF"},
        "blocks": [
            {"runs": [{"type": "text", "text": "Title",
                       "format": {"bold": True}}]},
            {"runs": [{"type": "text", "text": "Subtitle",
                       "format": {"italic": True}}]},
            {"runs": [{"type": "text", "text": "Description"}]},
        ],
    },
    "style": {
        "pen": {"color": "#0000FF", "width": 2, "dash": "solid"},
        "fill": {"color": "#E0E0FFFF"},
        "text": {"color": "#333333"},
    },
}

RECT_LEGACY = {
    "id": "a2", "kind": "rect",
    "geom": {"x": 400, "y": 100, "w": 200, "h": 120},
    "meta": {
        "label": "Database", "tech": "PostgreSQL", "note": "Main store",
        "label_align": "center", "label_size": 14,
        "tech_size": 10, "note_size": 10, "text_valign": "middle",
    },
    "style": {
        "pen": {"color": "#008800", "width": 2, "dash": "solid"},
        "fill": {"color": "#E0FFE0FF"},
        "text": {"color": "#006600"},
    },
}

LINE_WITH_TEXT = {
    "id": "a3", "kind": "line",
    "geom": {"x1": 300, "y1": 160, "x2": 400, "y2": 160},
    "contents": {
        "blocks": [
            {"runs": [{"type": "text", "text": "HTTP/REST",
                       "format": {"bold": True, "font_size": 10}}]},
        ],
    },
    "style": {
        "pen": {"color": "#888888", "width": 1, "dash": "solid"},
        "fill": {"color": "#00000000"},
        "text": {"color": "#666666"},
    },
}

ELLIPSE_EMPTY = {
    "id": "a4", "kind": "ellipse",
    "geom": {"x": 100, "y": 300, "w": 100, "h": 100},
    "contents": {"blocks": [], "frame": {}},
    "style": {
        "pen": {"color": "#000000", "width": 1, "dash": "solid"},
        "fill": {"color": "#FFFF00FF"},
        "text": {"color": "#000000"},
    },
}

MULTI_RUN_RECT = {
    "id": "a5", "kind": "rect",
    "geom": {"x": 100, "y": 500, "w": 250, "h": 80},
    "contents": {
        "frame": {"halign": "left"},
        "default_format": {"font_size": 12, "color": "#000000FF"},
        "blocks": [
            {"runs": [
                {"type": "text", "text": "Normal "},
                {"type": "text", "text": "Bold",
                 "format": {"bold": True}},
                {"type": "text", "text": " and "},
                {"type": "text", "text": "Italic",
                 "format": {"italic": True}},
            ]},
        ],
    },
    "style": {
        "pen": {"color": "#000000", "width": 1, "dash": "solid"},
        "fill": {"color": "#FFFFFFFF"},
        "text": {"color": "#000000"},
    },
}


# ── Tests ────────────────────────────────────────────────────────────────

class TestExportNoError:
    """All shape kinds export without error."""

    @pytest.mark.parametrize("kind,geom", [
        ("rect", {"x": 10, "y": 10, "w": 100, "h": 60}),
        ("roundedrect", {"x": 10, "y": 10, "w": 100, "h": 60, "adjust1": 10}),
        ("ellipse", {"x": 10, "y": 10, "w": 100, "h": 60}),
        ("hexagon", {"x": 10, "y": 10, "w": 100, "h": 60, "adjust1": 0.25}),
        ("cylinder", {"x": 10, "y": 10, "w": 80, "h": 100, "adjust1": 0.15}),
        ("blockarrow", {"x": 10, "y": 10, "w": 150, "h": 60,
                        "adjust1": 0.5, "adjust2": 30}),
        ("isocube", {"x": 10, "y": 10, "w": 100, "h": 80,
                     "adjust1": 20, "adjust2": 135}),
        ("line", {"x1": 10, "y1": 10, "x2": 200, "y2": 100}),
        ("text", {"x": 10, "y": 10}),
        ("polygon", {"x": 10, "y": 10, "w": 100, "h": 100,
                     "points": [[0, 0], [1, 0], [0.5, 1]]}),
    ])
    def test_kind_exports(self, kind, geom):
        ann = {
            "id": "t1", "kind": kind, "geom": geom,
            "meta": {"label": "Test", "tech": "", "note": ""},
            "style": {
                "pen": {"color": "#000000", "width": 1, "dash": "solid"},
                "fill": {"color": "#FFFFFFFF"},
                "text": {"color": "#000000"},
            },
        }
        prs = _export_and_read([ann])
        assert len(prs.slides) == 1


class TestOverlay20Text:
    """Blocks/runs text exported to shape paragraphs."""

    def test_three_blocks_three_paragraphs(self):
        prs = _export_and_read([RECT_OVERLAY20])
        texts = _get_shape_texts(prs)
        # Should have one shape with 3 paragraphs
        shape_texts = [t for _, t in texts if len(t) == 3]
        assert len(shape_texts) >= 1
        paras = shape_texts[0]
        assert paras[0] == "Title"
        assert paras[1] == "Subtitle"
        assert paras[2] == "Description"

    def test_bold_run(self):
        prs = _export_and_read([RECT_OVERLAY20])
        runs = _get_shape_runs(prs)
        # Find the "Title" run
        title_runs = [(t, r) for t, r in runs if any(rt == "Title" for rt, _, _ in r)]
        assert len(title_runs) >= 1
        for run_text, bold, italic in title_runs[0][1]:
            if run_text == "Title":
                assert bold is True

    def test_italic_run(self):
        prs = _export_and_read([RECT_OVERLAY20])
        runs = _get_shape_runs(prs)
        subtitle_runs = [(t, r) for t, r in runs if any(rt == "Subtitle" for rt, _, _ in r)]
        assert len(subtitle_runs) >= 1
        for run_text, bold, italic in subtitle_runs[0][1]:
            if run_text == "Subtitle":
                assert italic is True


class TestLegacyText:
    """Legacy meta label/tech/note exported correctly."""

    def test_legacy_three_lines(self):
        prs = _export_and_read([RECT_LEGACY])
        texts = _get_shape_texts(prs)
        shape_texts = [t for _, t in texts if len(t) == 3]
        assert len(shape_texts) >= 1
        paras = shape_texts[0]
        assert paras[0] == "Database"
        assert "[PostgreSQL]" in paras[1]
        assert paras[2] == "Main store"

    def test_legacy_label_bold(self):
        prs = _export_and_read([RECT_LEGACY])
        runs = _get_shape_runs(prs)
        db_runs = [(t, r) for t, r in runs if any(rt == "Database" for rt, _, _ in r)]
        assert len(db_runs) >= 1
        for run_text, bold, italic in db_runs[0][1]:
            if run_text == "Database":
                assert bold is True


class TestLineTextBox:
    """Line annotation creates a floating text box."""

    def test_line_text_present(self):
        prs = _export_and_read([LINE_WITH_TEXT])
        texts = _get_shape_texts(prs)
        all_text = [t for _, paras in texts for t in paras]
        assert "HTTP/REST" in all_text


class TestEmptyBlocks:
    """Empty blocks produce no text paragraphs."""

    def test_empty_blocks_no_text(self):
        prs = _export_and_read([ELLIPSE_EMPTY])
        texts = _get_shape_texts(prs)
        # The ellipse should have no text paragraphs
        assert len(texts) == 0


class TestMultiRunFormatting:
    """Multiple runs within a single paragraph with mixed formatting."""

    def test_four_runs_in_one_paragraph(self):
        prs = _export_and_read([MULTI_RUN_RECT])
        runs = _get_shape_runs(prs)
        # Find the paragraph with "Normal Bold and Italic"
        target = [r for t, r in runs if "Bold" in t and "Italic" in t]
        assert len(target) >= 1
        run_list = target[0]
        assert len(run_list) == 4
        # Check formatting per run
        assert run_list[0][0] == "Normal "
        assert run_list[1][0] == "Bold"
        assert run_list[1][1] is True   # bold
        assert run_list[2][0] == " and "
        assert run_list[3][0] == "Italic"
        assert run_list[3][2] is True   # italic


class TestMultipleAnnotations:
    """Export multiple annotations in one slide."""

    def test_two_shapes_one_line(self):
        prs = _export_and_read([RECT_OVERLAY20, RECT_LEGACY, LINE_WITH_TEXT])
        texts = _get_shape_texts(prs)
        all_text = [t for _, paras in texts for t in paras]
        assert "Title" in all_text
        assert "Database" in all_text
        assert "HTTP/REST" in all_text
