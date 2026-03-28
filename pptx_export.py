"""
pptx_export.py

Export canvas annotations to PowerPoint slides.
"""

from __future__ import annotations

from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsmap


# PowerPoint uses EMUs (English Metric Units)
# 914400 EMUs = 1 inch, typical screen is 96 DPI
# So 1 pixel at 96 DPI = 914400 / 96 = 9525 EMUs
EMUS_PER_PIXEL = 9525

# Default slide size (standard 16:9 widescreen)
DEFAULT_SLIDE_WIDTH = Emu(12192000)  # 10 inches
DEFAULT_SLIDE_HEIGHT = Emu(6858000)  # 7.5 inches


def hex_to_rgb(hex_color: str) -> Optional[Tuple[int, int, int]]:
    """
    Convert hex color string to RGB tuple.

    Args:
        hex_color: Color in format "#RRGGBB" or "#RRGGBBAA"

    Returns:
        Tuple of (r, g, b) or None if invalid
    """
    if not hex_color or not hex_color.startswith("#"):
        return None

    hex_color = hex_color.lstrip("#")

    # Handle alpha channel (#RRGGBBAA format) — strip trailing alpha
    if len(hex_color) == 8:
        hex_color = hex_color[:6]

    if len(hex_color) != 6:
        return None

    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return (r, g, b)
    except ValueError:
        return None


def get_alpha(hex_color: str) -> int:
    """
    Return the alpha byte (0–255) from a hex color string.

    Args:
        hex_color: Color in format "#RRGGBB" or "#RRGGBBAA"

    Returns:
        Alpha value 0–255 (255 = fully opaque, default when no alpha present)
    """
    if not hex_color or not hex_color.startswith("#"):
        return 255

    raw = hex_color.lstrip("#")

    if len(raw) == 8:
        try:
            return int(raw[6:8], 16)
        except ValueError:
            return 255

    return 255


def _get_contents(record: Dict[str, Any]) -> Dict[str, Any]:
    """Extract the contents/meta dict from an annotation record.

    Handles both overlay-2.0 (``contents`` key with ``blocks``) and
    legacy (``meta`` key with flat ``label``/``tech``/``note`` fields).
    """
    return record.get("contents") or record.get("meta") or {}


def _get_blocks(record: Dict[str, Any]) -> list:
    """Extract blocks list from an annotation record.

    Returns the ``blocks`` array from overlay-2.0 ``contents`` dict,
    or builds a synthetic 3-block list from legacy ``meta`` fields.

    Each returned block is a dict: ``{"runs": [...], "halign": "..."}``
    """
    contents = _get_contents(record)
    if "blocks" in contents and contents["blocks"]:
        return contents["blocks"]
    # Build synthetic blocks from legacy fields
    blocks = []
    label = contents.get("label", "")
    tech = contents.get("tech", "")
    note = contents.get("note", "")
    if label:
        blocks.append({"runs": [{"type": "text", "text": label,
                                  "format": {"bold": True}}]})
    if tech:
        blocks.append({"runs": [{"type": "text", "text": f"[{tech}]",
                                  "format": {"italic": True}}]})
    if note:
        blocks.append({"runs": [{"type": "text", "text": note}]})
    return blocks


def _get_default_format(record: Dict[str, Any]) -> Dict[str, Any]:
    """Get the default_format from contents, or build from legacy/style."""
    contents = _get_contents(record)
    df = contents.get("default_format", {})
    if df:
        return df
    style = record.get("style", {})
    text_style = style.get("text", {})
    return {
        "font_family": contents.get("font_family", ""),
        "font_size": contents.get("font_size", 12),
        "color": contents.get("color", text_style.get("color", "#000000")),
    }


def _get_frame(record: Dict[str, Any]) -> Dict[str, Any]:
    """Get the frame from contents, or build from legacy fields."""
    contents = _get_contents(record)
    if "frame" in contents and contents["frame"]:
        return contents["frame"]
    return {
        "halign": contents.get("halign", contents.get("label_align", "center")),
        "valign": contents.get("valign", contents.get("text_valign", "top")),
    }


def has_alpha_transparency(hex_color: str) -> bool:
    """
    Check if a hex color has alpha transparency (not fully opaque).

    Args:
        hex_color: Color in format "#RRGGBB" or "#RRGGBBAA"

    Returns:
        True if the color has transparency (alpha < FF)
    """
    return get_alpha(hex_color) < 255


def px_to_emu(pixels: float) -> int:
    """Convert pixels to EMUs."""
    return int(pixels * EMUS_PER_PIXEL)


def _apply_line_style(line, style: Dict[str, Any], arrow_mode: str = "none"):
    """
    Apply stroke/line styling to a shape's line.

    Args:
        line: The shape.line object from python-pptx
        style: Style dict containing pen properties
        arrow_mode: Arrow mode for lines (none, start, end, both)
    """
    pen = style.get("pen", {})

    # Set line color
    pen_color = pen.get("color", "#FF0000")
    rgb = hex_to_rgb(pen_color)
    if rgb:
        line.color.rgb = RGBColor(*rgb)

    # Set line width
    pen_width = pen.get("width", 2)
    line.width = Pt(pen_width)

    # Set dash style
    dash = pen.get("dash", "solid")
    if dash == "dashed":
        # Use round dot dash style for dashed lines
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        line.dash_style = MSO_LINE_DASH_STYLE.DASH


def _apply_fill_style(fill, style: Dict[str, Any]):
    """
    Apply fill styling to a shape.

    Handles three cases:
      - Fully transparent (alpha == 0): no fill
      - Semi-transparent (0 < alpha < 255): solid fill with PPTX transparency
      - Fully opaque (alpha == 255 or no alpha): solid fill

    Args:
        fill: The shape.fill object from python-pptx
        style: Style dict containing fill properties
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    fill_style = style.get("fill") or {}
    fill_color = fill_style.get("color", "#00000000")

    alpha = get_alpha(fill_color)

    if alpha == 0:
        fill.background()  # Fully transparent — no fill
        return

    rgb = hex_to_rgb(fill_color)
    if not rgb:
        fill.background()
        return

    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)

    # Apply transparency for semi-transparent fills
    if alpha < 255:
        # PPTX transparency is expressed as alpha in 1/1000 % (0 = opaque, 100000 = invisible)
        pptx_alpha = int(round(alpha / 255 * 100000))
        # Set a:alpha on the srgbClr element
        srgb_el = fill._fill.find(f'.//{qn("a:srgbClr")}')
        if srgb_el is not None:
            # Remove any existing alpha child
            for old in srgb_el.findall(qn("a:alpha")):
                srgb_el.remove(old)
            alpha_el = etree.SubElement(srgb_el, qn("a:alpha"))
            alpha_el.set("val", str(pptx_alpha))


def _apply_rotation(shape, geom: Dict[str, Any]):
    """Apply rotation angle from geom dict to a PPTX shape.

    Args:
        shape: A python-pptx shape object.
        geom: The annotation's geom dict (may contain 'angle' in degrees).
    """
    angle = geom.get("angle", 0)
    if angle:
        shape.rotation = float(angle)


def _add_text_to_shape(shape, record: Dict[str, Any], text_style: Dict[str, Any]):
    """Add text content from blocks/runs to a shape's text frame.

    Reads the overlay-2.0 ``contents.blocks`` structure or falls back
    to legacy ``meta.label``/``tech``/``note`` fields.  Each block maps
    to a PowerPoint paragraph; each run within a block maps to a
    PowerPoint text run with per-run formatting (bold, italic, color,
    font family, font size).

    Args:
        shape: The PowerPoint shape (must have text_frame).
        record: The full annotation record dict.
        text_style: Fallback text style from ``style.text``.
    """
    if not shape.has_text_frame:
        return

    blocks = _get_blocks(record)
    if not blocks:
        return

    df = _get_default_format(record)
    frame = _get_frame(record)

    tf = shape.text_frame
    tf.word_wrap = True

    # Vertical alignment
    valign = frame.get("valign", "top")
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    # Default text properties
    def_color = df.get("color", text_style.get("color", "#000000"))
    def_rgb = hex_to_rgb(def_color)
    def_size = df.get("font_size", 12)
    def_family = df.get("font_family", "")
    def_halign = frame.get("halign", "center")

    _align_map = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT,
                  "center": PP_ALIGN.CENTER, "justified": PP_ALIGN.JUSTIFY}

    first_para = True
    for blk in blocks:
        runs_data = blk.get("runs", []) if isinstance(blk, dict) else getattr(blk, "runs", [])
        # Skip empty blocks
        if not runs_data:
            continue
        has_text = any(
            (r.get("text", "") if isinstance(r, dict) else getattr(r, "text", ""))
            for r in runs_data
            if (r.get("type", "text") if isinstance(r, dict) else getattr(r, "type", "text")) == "text"
        )
        if not has_text:
            continue

        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()
        p.clear()

        # Block-level alignment
        blk_halign = (blk.get("halign", "") if isinstance(blk, dict)
                      else getattr(blk, "halign", "")) or def_halign
        p.alignment = _align_map.get(blk_halign, PP_ALIGN.CENTER)

        # Block-level spacing
        sb = (blk.get("space_before", 0) if isinstance(blk, dict)
              else getattr(blk, "space_before", 0))
        sa = (blk.get("space_after", 0) if isinstance(blk, dict)
              else getattr(blk, "space_after", 0))
        if sb:
            p.space_before = Pt(int(sb))
        if sa:
            p.space_after = Pt(int(sa))

        # Add runs
        for run_data in runs_data:
            if isinstance(run_data, dict):
                rtype = run_data.get("type", "text")
                rtext = run_data.get("text", "")
                rfmt = run_data.get("format", {}) or {}
            else:
                rtype = getattr(run_data, "type", "text")
                rtext = getattr(run_data, "text", "")
                rf = getattr(run_data, "format", None)
                rfmt = rf.to_dict() if rf and hasattr(rf, "to_dict") else {}

            if rtype != "text" or not rtext:
                continue

            run = p.add_run()
            run.text = rtext

            # Per-run formatting (sparse overrides of default_format)
            run_size = rfmt.get("font_size", def_size)
            run.font.size = Pt(int(run_size))

            run_bold = rfmt.get("bold", False)
            if run_bold:
                run.font.bold = True

            run_italic = rfmt.get("italic", False)
            if run_italic:
                run.font.italic = True

            run_underline = rfmt.get("underline", False)
            if run_underline:
                run.font.underline = True

            run_color = rfmt.get("color", "")
            run_rgb = hex_to_rgb(run_color) if run_color else def_rgb
            if run_rgb:
                run.font.color.rgb = RGBColor(*run_rgb)

            run_family = rfmt.get("font_family", def_family)
            if run_family:
                run.font.name = run_family


def _add_rectangle(slide, record: Dict[str, Any], scale_x: float, scale_y: float):
    """Add a rectangle shape to the slide."""
    geom = record.get("geom", {})
    x = px_to_emu(geom.get("x", 0) * scale_x)
    y = px_to_emu(geom.get("y", 0) * scale_y)
    w = px_to_emu(geom.get("w", 100) * scale_x)
    h = px_to_emu(geom.get("h", 100) * scale_y)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)

    style = record.get("style", {})
    _apply_line_style(shape.line, style)
    _apply_fill_style(shape.fill, style)

    _add_text_to_shape(shape, record, style.get("text", {}))
    _apply_rotation(shape, geom)


def _add_rounded_rectangle(slide, record: Dict[str, Any], scale_x: float, scale_y: float):
    """Add a rounded rectangle shape to the slide."""
    geom = record.get("geom", {})
    x = px_to_emu(geom.get("x", 0) * scale_x)
    y = px_to_emu(geom.get("y", 0) * scale_y)
    w = px_to_emu(geom.get("w", 100) * scale_x)
    h = px_to_emu(geom.get("h", 100) * scale_y)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)

    # Adjust corner radius if possible
    # python-pptx doesn't directly expose corner radius, but we can try via adjustments
    radius = geom.get("adjust1", geom.get("radius", 10))  # Support legacy "radius" key
    try:
        # The adjustment value is a percentage (0-100000)
        # Corner radius is relative to the shorter dimension
        min_dim = min(geom.get("w", 100), geom.get("h", 100))
        if min_dim > 0:
            adj_value = int((radius / min_dim) * 100000)
            adj_value = min(50000, max(0, adj_value))  # Clamp to valid range
            shape.adjustments[0] = adj_value / 100000
    except (IndexError, AttributeError):
        pass  # Not all shapes support adjustments

    style = record.get("style", {})
    _apply_line_style(shape.line, style)
    _apply_fill_style(shape.fill, style)

    _add_text_to_shape(shape, record, style.get("text", {}))
    _apply_rotation(shape, geom)


def _add_ellipse(slide, record: Dict[str, Any], scale_x: float, scale_y: float):
    """Add an ellipse/oval shape to the slide."""
    geom = record.get("geom", {})
    x = px_to_emu(geom.get("x", 0) * scale_x)
    y = px_to_emu(geom.get("y", 0) * scale_y)
    w = px_to_emu(geom.get("w", 100) * scale_x)
    h = px_to_emu(geom.get("h", 100) * scale_y)

    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)

    style = record.get("style", {})
    _apply_line_style(shape.line, style)
    _apply_fill_style(shape.fill, style)

    _add_text_to_shape(shape, record, style.get("text", {}))
    _apply_rotation(shape, geom)


def _add_line(slide, record: Dict[str, Any], scale_x: float, scale_y: float):
    """Add a line connector to the slide."""
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    from lxml import etree

    geom = record.get("geom", {})
    x1 = px_to_emu(geom.get("x1", 0) * scale_x)
    y1 = px_to_emu(geom.get("y1", 0) * scale_y)
    x2 = px_to_emu(geom.get("x2", 100) * scale_x)
    y2 = px_to_emu(geom.get("y2", 100) * scale_y)

    # add_connector expects (connector_type, begin_x, begin_y, end_x, end_y)
    # Pass the actual endpoints directly to preserve line direction
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        x1, y1,
        x2, y2
    )

    # Apply line style
    style = record.get("style", {})
    _apply_line_style(connector.line, style)

    # Handle arrows using XML manipulation (python-pptx doesn't expose this directly)
    arrow_mode = style.get("arrow", "none")
    if arrow_mode != "none":
        # Access the line element in the connector's spPr
        spPr = connector._element.spPr
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            # Create ln element if it doesn't exist
            ln = etree.SubElement(spPr, qn("a:ln"))

        # Add head (start) arrow — must come before tailEnd in OOXML
        if arrow_mode in ("start", "both"):
            head_end = ln.find(qn("a:headEnd"))
            if head_end is None:
                head_end = etree.SubElement(ln, qn("a:headEnd"))
            head_end.set("type", "triangle")

        # Add tail (end) arrow
        if arrow_mode in ("end", "both"):
            tail_end = ln.find(qn("a:tailEnd"))
            if tail_end is None:
                tail_end = etree.SubElement(ln, qn("a:tailEnd"))
            tail_end.set("type", "triangle")

    # Add text label as a text box at the midpoint of the line
    _add_line_curve_textbox(slide, record, geom, scale_x, scale_y)


def _add_line_curve_textbox(slide, record: Dict[str, Any],
                            geom: Dict[str, Any],
                            scale_x: float, scale_y: float,
                            mid_x: float | None = None,
                            mid_y: float | None = None):
    """Add a floating text box for a line or curve annotation.

    Creates a textbox at the midpoint of the line/curve and populates it
    from ``contents.blocks.runs`` (or legacy ``meta.label``/``tech``/``note``).

    Args:
        slide: The PowerPoint slide.
        record: Full annotation record.
        geom: Geometry dict (for line: x1/y1/x2/y2; for curve: x/y/w/h).
        scale_x, scale_y: Pixel-to-slide scale factors.
        mid_x, mid_y: Pre-computed midpoint in pixels (for curves).
                       If None, computed from line endpoints.
    """
    blocks = _get_blocks(record)
    if not blocks:
        return

    df = _get_default_format(record)
    def_size = df.get("font_size", 12)
    def_color = df.get("color", "#000000")
    def_rgb = hex_to_rgb(def_color)

    # Compute midpoint if not provided
    if mid_x is None or mid_y is None:
        mid_x = (geom.get("x1", 0) + geom.get("x2", 0)) / 2
        mid_y = (geom.get("y1", 0) + geom.get("y2", 0)) / 2

    # Estimate text box dimensions
    contents = _get_contents(record)
    box_w = contents.get("text_box_width", 0)
    box_h = contents.get("text_box_height", 0)
    if not box_w or not box_h:
        # Estimate from block text
        all_text = ""
        n_lines = 0
        for blk in blocks:
            runs = blk.get("runs", []) if isinstance(blk, dict) else getattr(blk, "runs", [])
            line_text = ""
            for r in runs:
                t = r.get("text", "") if isinstance(r, dict) else getattr(r, "text", "")
                line_text += t
            if line_text:
                n_lines += 1
                if len(line_text) > len(all_text):
                    all_text = line_text
        if not box_w:
            box_w = max(50, len(all_text) * def_size * 0.6 + 20)
        if not box_h:
            box_h = n_lines * def_size * 1.5 + 10

    tx = px_to_emu((mid_x - box_w / 2) * scale_x)
    ty = px_to_emu((mid_y - box_h / 2) * scale_y)
    tw = px_to_emu(box_w * scale_x)
    th = px_to_emu(box_h * scale_y)

    tbox = slide.shapes.add_textbox(tx, ty, tw, th)
    # Reuse the shape text renderer
    _add_text_to_shape(tbox, record, record.get("style", {}).get("text", {}))


def _add_text(slide, record: Dict[str, Any], scale_x: float, scale_y: float):
    """Add a text box to the slide."""
    geom = record.get("geom", {})
    x = px_to_emu(geom.get("x", 0) * scale_x)
    y = px_to_emu(geom.get("y", 0) * scale_y)

    text = record.get("meta", {}).get("note", "") or record.get("text", "")
    style = record.get("style", {})
    text_style = style.get("text", {})
    font_size = text_style.get("size_pt", 12)

    # Use geom dimensions when available, otherwise estimate
    if geom.get("w"):
        w = px_to_emu(geom["w"] * scale_x)
    else:
        char_width = font_size * 0.6
        w = px_to_emu(max(50, len(text) * char_width + 20) * scale_x)

    if geom.get("h"):
        h = px_to_emu(geom["h"] * scale_y)
    else:
        h = px_to_emu((font_size * 1.5 + 10) * scale_y)

    shape = slide.shapes.add_textbox(x, y, w, h)

    # Text items have no outline border
    shape.line.fill.background()

    # Apply fill style (usually transparent for text items)
    _apply_fill_style(shape.fill, style)

    tf = shape.text_frame
    tf.word_wrap = True

    # Vertical alignment (python-pptx 1.0.2 uses vertical_anchor, not anchor)
    meta = record.get("meta", {})
    valign = meta.get("text_valign", "top")
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.clear()
    run = p.add_run()
    run.text = text

    # Apply text styling
    text_color = text_style.get("color", "#000000")
    rgb = hex_to_rgb(text_color)
    run.font.size = Pt(font_size)
    if rgb:
        run.font.color.rgb = RGBColor(*rgb)

    # Horizontal alignment
    align = meta.get("label_align", "center")
    if align == "left":
        p.alignment = PP_ALIGN.LEFT
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.CENTER


def _add_hexagon(slide, record: Dict[str, Any], scale_x: float, scale_y: float):
    """Add a hexagon shape to the slide."""
    geom = record.get("geom", {})
    x = px_to_emu(geom.get("x", 0) * scale_x)
    y = px_to_emu(geom.get("y", 0) * scale_y)
    w = px_to_emu(geom.get("w", 100) * scale_x)
    h = px_to_emu(geom.get("h", 80) * scale_y)

    shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, x, y, w, h)

    # Adjust indent ratio (adjust1 is 0.0-0.5, PPTX adjustment is 0.0-0.5)
    adjust1 = geom.get("adjust1", 0.25)
    try:
        shape.adjustments[0] = adjust1
    except (IndexError, AttributeError):
        pass

    style = record.get("style", {})
    _apply_line_style(shape.line, style)
    _apply_fill_style(shape.fill, style)

    _add_text_to_shape(shape, record, style.get("text", {}))
    _apply_rotation(shape, geom)


def _add_cylinder(slide, record: Dict[str, Any], scale_x: float, scale_y: float):
    """Add a cylinder (can) shape to the slide."""
    geom = record.get("geom", {})
    x = px_to_emu(geom.get("x", 0) * scale_x)
    y = px_to_emu(geom.get("y", 0) * scale_y)
    w = px_to_emu(geom.get("w", 80) * scale_x)
    h = px_to_emu(geom.get("h", 120) * scale_y)

    shape = slide.shapes.add_shape(MSO_SHAPE.CAN, x, y, w, h)

    # Adjust cap height ratio (adjust1 is 0.1-0.5, PPTX adjustment is similar)
    adjust1 = geom.get("adjust1", 0.15)
    try:
        shape.adjustments[0] = adjust1
    except (IndexError, AttributeError):
        pass

    style = record.get("style", {})
    _apply_line_style(shape.line, style)
    _apply_fill_style(shape.fill, style)

    _add_text_to_shape(shape, record, style.get("text", {}))
    _apply_rotation(shape, geom)


def _add_blockarrow(slide, record: Dict[str, Any], scale_x: float, scale_y: float):
    """Add a block arrow shape to the slide."""
    geom = record.get("geom", {})
    x = px_to_emu(geom.get("x", 0) * scale_x)
    y = px_to_emu(geom.get("y", 0) * scale_y)
    w = px_to_emu(geom.get("w", 150) * scale_x)
    h = px_to_emu(geom.get("h", 60) * scale_y)

    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)

    # Adjust shaft width ratio and arrowhead length
    adjust1 = geom.get("adjust1", 0.5)  # shaft width ratio
    adjust2 = geom.get("adjust2", 15)   # arrowhead length in pixels
    try:
        # PPTX RIGHT_ARROW: adj1 = shaft width ratio, adj2 = head length ratio
        shape.adjustments[0] = adjust1
        total_w = geom.get("w", 150)
        if total_w > 0:
            shape.adjustments[1] = adjust2 / total_w
    except (IndexError, AttributeError):
        pass

    style = record.get("style", {})
    _apply_line_style(shape.line, style)
    _apply_fill_style(shape.fill, style)

    _add_text_to_shape(shape, record, style.get("text", {}))
    _apply_rotation(shape, geom)


def _add_polygon(slide, record: Dict[str, Any], scale_x: float, scale_y: float):
    """Add a polygon (freeform) shape to the slide."""
    geom = record.get("geom", {})
    x = geom.get("x", 0) * scale_x
    y = geom.get("y", 0) * scale_y
    w = geom.get("w", 100) * scale_x
    h = geom.get("h", 100) * scale_y

    points = geom.get("points", [])
    if not points or len(points) < 3:
        # Fall back to rectangle if no valid polygon points
        _add_rectangle(slide, record, scale_x, scale_y)
        return

    # Build freeform shape from relative points
    # points are [[rx, ry], ...] where rx/ry are 0.0-1.0 relative to w/h
    start_x = px_to_emu(x + points[0][0] * w)
    start_y = px_to_emu(y + points[0][1] * h)
    builder = slide.shapes.build_freeform(start_x, start_y)

    # add_line_segments takes an iterable of (x, y) pairs and auto-closes
    vertices = [
        (px_to_emu(x + pt[0] * w), px_to_emu(y + pt[1] * h))
        for pt in points[1:]
    ]
    builder.add_line_segments(vertices, close=True)

    shape = builder.convert_to_shape()

    style = record.get("style", {})
    _apply_line_style(shape.line, style)
    _apply_fill_style(shape.fill, style)

    _add_text_to_shape(shape, record, style.get("text", {}))
    _apply_rotation(shape, geom)


def _add_curve(slide, record: Dict[str, Any], scale_x: float, scale_y: float):
    """Add a curve (freeform path) shape to the slide.

    Emits native OOXML ``a:cubicBezTo`` and ``a:quadBezTo`` elements so
    that cubic and quadratic bezier control points are preserved in the
    PowerPoint file.
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    geom = record.get("geom", {})
    x = geom.get("x", 0) * scale_x
    y = geom.get("y", 0) * scale_y
    w = geom.get("w", 100) * scale_x
    h = geom.get("h", 100) * scale_y

    nodes = geom.get("nodes", [])
    if not nodes or len(nodes) < 2:
        return

    # Collect ALL coordinates (anchors + control points) so the freeform
    # builder computes a bounding box that encompasses everything.
    first_node = nodes[0]
    fx = first_node.get("x", 0)
    fy = first_node.get("y", 0)
    start_x = px_to_emu(x + fx * w)
    start_y = px_to_emu(y + fy * h)
    builder = slide.shapes.build_freeform(start_x, start_y)

    all_vertices: List[Tuple[float, float]] = []
    for node in nodes[1:]:
        if node.get("cmd") == "Z":
            continue
        all_vertices.append(
            (px_to_emu(x + node.get("x", 0) * w),
             px_to_emu(y + node.get("y", 0) * h))
        )
        for kx, ky in (("c1x", "c1y"), ("c2x", "c2y"), ("cx", "cy")):
            if kx in node:
                all_vertices.append(
                    (px_to_emu(x + node[kx] * w),
                     px_to_emu(y + node[ky] * h))
                )

    if not all_vertices:
        return

    builder.add_line_segments(all_vertices, close=False)
    shape = builder.convert_to_shape()

    # ── Replace the builder's line-only path with proper bezier XML ──
    sp = shape._element
    path_el = sp.find(f'.//{qn("a:path")}')
    if path_el is not None:
        # Read shape offset so path coordinates are relative to shape origin
        offset_x = int(shape.left)
        offset_y = int(shape.top)

        def _px(rel: float) -> str:
            return str(int(round(px_to_emu(x + rel * w) - offset_x)))

        def _py(rel: float) -> str:
            return str(int(round(px_to_emu(y + rel * h) - offset_y)))

        def _pt(parent: etree._Element, rx: float, ry: float) -> None:
            pt = etree.SubElement(parent, qn("a:pt"))
            pt.set("x", _px(rx))
            pt.set("y", _py(ry))

        # Clear existing path children (moveTo + lnTo from builder)
        for child in list(path_el):
            path_el.remove(child)

        # Rebuild with native bezier elements
        move = etree.SubElement(path_el, qn("a:moveTo"))
        _pt(move, fx, fy)

        for node in nodes[1:]:
            cmd = node.get("cmd", "L")

            if cmd == "C":
                el = etree.SubElement(path_el, qn("a:cubicBezTo"))
                _pt(el, node["c1x"], node["c1y"])
                _pt(el, node["c2x"], node["c2y"])
                _pt(el, node["x"], node["y"])

            elif cmd == "Q":
                el = etree.SubElement(path_el, qn("a:quadBezTo"))
                _pt(el, node["cx"], node["cy"])
                _pt(el, node["x"], node["y"])

            elif cmd == "L":
                el = etree.SubElement(path_el, qn("a:lnTo"))
                _pt(el, node["x"], node["y"])

            elif cmd == "M":
                el = etree.SubElement(path_el, qn("a:moveTo"))
                _pt(el, node["x"], node["y"])

            elif cmd == "Z":
                etree.SubElement(path_el, qn("a:close"))

    _apply_rotation(shape, geom)

    style = record.get("style", {})
    _apply_line_style(shape.line, style)
    # No fill for curves
    shape.fill.background()

    # Arrowheads — same XML approach used by _add_line
    arrow_mode = style.get("arrow", "none")
    if arrow_mode != "none":
        spPr = sp.find(qn("a:spPr")) if sp.find(qn("a:spPr")) is not None else sp.spPr
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = etree.SubElement(spPr, qn("a:ln"))
        # headEnd must come before tailEnd in OOXML element order
        if arrow_mode in ("start", "both"):
            head_end = ln.find(qn("a:headEnd"))
            if head_end is None:
                head_end = etree.SubElement(ln, qn("a:headEnd"))
            head_end.set("type", "triangle")
        if arrow_mode in ("end", "both"):
            tail_end = ln.find(qn("a:tailEnd"))
            if tail_end is None:
                tail_end = etree.SubElement(ln, qn("a:tailEnd"))
            tail_end.set("type", "triangle")

    # ── Text label at curve midpoint ──
    # Evaluate the actual curve midpoint (halfway along the path)
    geom = record.get("geom", {})
    gx = geom.get("x", 0)
    gy = geom.get("y", 0)
    gw = geom.get("w", 0)
    gh = geom.get("h", 0)

    segments: List[Tuple[Tuple[float, float], Dict[str, Any]]] = []
    cur_x, cur_y = nodes[0].get("x", 0), nodes[0].get("y", 0)
    for nd in nodes[1:]:
        cmd = nd.get("cmd", "L")
        if cmd == "Z":
            continue
        segments.append(((cur_x, cur_y), nd))
        cur_x, cur_y = nd.get("x", cur_x), nd.get("y", cur_y)

    if segments:
        seg_start, seg_node = segments[len(segments) // 2]
        cmd = seg_node.get("cmd", "L")
        t = 0.5
        if cmd == "C":
            p0x, p0y = seg_start
            p1x, p1y = seg_node["c1x"], seg_node["c1y"]
            p2x, p2y = seg_node["c2x"], seg_node["c2y"]
            p3x, p3y = seg_node["x"], seg_node["y"]
            u = 1 - t
            rx = u**3*p0x + 3*u**2*t*p1x + 3*u*t**2*p2x + t**3*p3x
            ry = u**3*p0y + 3*u**2*t*p1y + 3*u*t**2*p2y + t**3*p3y
        elif cmd == "Q":
            p0x, p0y = seg_start
            p1x, p1y = seg_node["cx"], seg_node["cy"]
            p2x, p2y = seg_node["x"], seg_node["y"]
            u = 1 - t
            rx = u**2*p0x + 2*u*t*p1x + t**2*p2x
            ry = u**2*p0y + 2*u*t*p1y + t**2*p2y
        else:
            rx = (seg_start[0] + seg_node.get("x", 0)) / 2
            ry = (seg_start[1] + seg_node.get("y", 0)) / 2
        mid_x = gx + rx * gw
        mid_y = gy + ry * gh
    else:
        mid_x = gx + gw / 2
        mid_y = gy + gh / 2

    _add_line_curve_textbox(slide, record, geom, scale_x, scale_y,
                            mid_x=mid_x, mid_y=mid_y)


def _add_seqblock(slide, record: Dict[str, Any], scale_x: float, scale_y: float):
    """Add a UML sequence-diagram combined fragment (alt/loop/opt/par/…).

    Built as a grouped compound shape: outer rectangle, pentagon type-tab,
    dashed divider lines, and section text labels.

    Args:
        slide: The PowerPoint slide
        record: Annotation record dict
        scale_x: Horizontal scale factor
        scale_y: Vertical scale factor
    """
    from pptx.enum.dml import MSO_LINE_DASH_STYLE

    geom = record.get("geom", {})
    meta = record.get("contents") or record.get("meta") or {}
    style = record.get("style", {})

    gx = geom.get("x", 0) * scale_x
    gy = geom.get("y", 0) * scale_y
    gw = geom.get("w", 100) * scale_x
    gh = geom.get("h", 100) * scale_y

    # Determine block type from DSL metadata or label
    dsl = record.get("dsl") or {}
    if not dsl:
        # Legacy: dsl may be inside meta.extras["dsl"]
        dsl = meta.get("extras", {}).get("dsl", {}) if isinstance(meta.get("extras"), dict) else {}
    block_type = dsl.get("sequence", {}).get("block_type", meta.get("label", "alt"))

    # Collect divider positions (ratios 0–1)
    dividers: List[float] = []
    for key in ("adjust1", "adjust2", "adjust3"):
        val = geom.get(key)
        if val is not None:
            dividers.append(float(val))

    # Pen / fill colours
    pen_style = style.get("pen", {})
    pen_color_hex = pen_style.get("color", "#9370DB")
    pen_rgb = hex_to_rgb(pen_color_hex)
    fill_style = style.get("fill", {})
    fill_color_hex = fill_style.get("color", "#ECECFF")
    fill_rgb = hex_to_rgb(fill_color_hex)
    text_style = style.get("text", {})
    text_color_hex = text_style.get("color", "#333333")
    text_rgb = hex_to_rgb(text_color_hex)

    # Convert to EMU
    ex = px_to_emu(gx)
    ey = px_to_emu(gy)
    ew = px_to_emu(gw)
    eh = px_to_emu(gh)

    # ── Build grouped shape ──
    grp = slide.shapes.add_group_shape()

    # 1) Outer rectangle — border + semi-transparent fill
    outer = grp.shapes.add_shape(MSO_SHAPE.RECTANGLE, ex, ey, ew, eh)
    if fill_rgb:
        outer.fill.solid()
        outer.fill.fore_color.rgb = RGBColor(*fill_rgb)
        # Apply alpha if present in the hex colour
        fill_hex_raw = fill_color_hex.lstrip("#")
        if len(fill_hex_raw) == 8:
            alpha_byte = int(fill_hex_raw[6:8], 16)
            alpha_pct = int(alpha_byte / 255 * 100000)
            from lxml import etree
            from pptx.oxml.ns import qn
            sf_el = outer._element.find(f'.//{qn("a:solidFill")}')
            if sf_el is not None and len(sf_el):
                alpha_el = etree.SubElement(sf_el[0], qn("a:alpha"))
                alpha_el.set("val", str(alpha_pct))
    else:
        outer.fill.background()
    if pen_rgb:
        outer.line.color.rgb = RGBColor(*pen_rgb)
    pen_width = pen_style.get("width", 2)
    outer.line.width = Pt(pen_width)

    # 2) Type tab — pentagon-like shape in top-left
    tab_text = block_type.upper()
    tab_font_size = 8
    tab_padding = 4
    # Approximate tab width from text length
    tab_w_px = len(tab_text) * tab_font_size * 0.7 + 2 * tab_padding + 8
    tab_h_px = 20
    tab_w = px_to_emu(tab_w_px * scale_x)
    tab_h = px_to_emu(tab_h_px * scale_y)

    # Use a freeform pentagon for the tab
    notch = px_to_emu(6 * scale_y)  # bottom-right notch
    builder = grp.shapes.build_freeform(ex, ey)
    builder.add_line_segments([
        (ex + tab_w, ey),
        (ex + tab_w, ey + tab_h - notch),
        (ex + tab_w - notch, ey + tab_h),
        (ex, ey + tab_h),
        (ex, ey),
    ], close=True)
    tab_shape = builder.convert_to_shape()
    if pen_rgb:
        tab_shape.fill.solid()
        tab_shape.fill.fore_color.rgb = RGBColor(*pen_rgb)
        # Semi-transparent tab fill
        from lxml import etree
        from pptx.oxml.ns import qn
        solid_fill = tab_shape._element.find(f'.//{qn("a:solidFill")}')
        if solid_fill is not None and len(solid_fill):
            alpha_el = etree.SubElement(solid_fill[0], qn("a:alpha"))
            alpha_el.set("val", "60000")  # 60% opaque
    tab_shape.line.color.rgb = RGBColor(*(pen_rgb or (0x93, 0x70, 0xDB)))
    tab_shape.line.width = Pt(pen_width)

    # Tab label text
    tab_label = grp.shapes.add_textbox(
        ex + px_to_emu(tab_padding * scale_x),
        ey,
        tab_w - px_to_emu(tab_padding * 2 * scale_x),
        tab_h,
    )
    tf = tab_label.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].text = tab_text
    tf.paragraphs[0].font.size = Pt(tab_font_size)
    tf.paragraphs[0].font.bold = True
    if text_rgb:
        tf.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
    from pptx.enum.text import MSO_ANCHOR
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # 3) Dashed divider lines
    for div_ratio in dividers:
        div_y = ey + int(eh * div_ratio)
        conn = grp.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            ex, div_y,
            ex + ew, div_y,
        )
        if pen_rgb:
            conn.line.color.rgb = RGBColor(*pen_rgb)
        conn.line.width = Pt(1)
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    # 4) Section labels (from meta.tech, pipe-separated)
    tech = meta.get("tech", "")
    sections = [s.strip() for s in tech.split("|")] if tech else []

    # Region top-Y values: section 0 starts below tab, others below dividers
    region_tops = [ey + tab_h + px_to_emu(4 * scale_y)]
    for div_ratio in dividers:
        region_tops.append(ey + int(eh * div_ratio) + px_to_emu(4 * scale_y))

    tech_size = meta.get("tech_size", 10)
    for i, sec_text in enumerate(sections):
        if not sec_text or i >= len(region_tops):
            continue
        sec_y = region_tops[i]
        sec_h = px_to_emu(tech_size * 2 * scale_y)
        sec_box = grp.shapes.add_textbox(
            ex + px_to_emu(4 * scale_x),
            sec_y,
            ew - px_to_emu(8 * scale_x),
            sec_h,
        )
        stf = sec_box.text_frame
        stf.word_wrap = True
        stf.paragraphs[0].text = f"[{sec_text}]"
        stf.paragraphs[0].font.size = Pt(tech_size)
        stf.paragraphs[0].font.italic = True
        if text_rgb:
            stf.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
        stf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 5) Set group extents and apply rotation
    grp.left = ex
    grp.top = ey
    grp.width = ew
    grp.height = eh
    _apply_rotation(grp, geom)


def _add_isocube(slide, record: Dict[str, Any], scale_x: float, scale_y: float):
    """Add an isometric cube shape to the slide.

    Maps to PowerPoint's CUBE auto-shape with flip attributes to approximate
    the extrusion angle.  The depth adjustment is mapped from our pixel depth
    to the PPTX proportional ratio (0.0–0.5).

    Args:
        slide: The PowerPoint slide
        record: Annotation record dict
        scale_x: Horizontal scale factor
        scale_y: Vertical scale factor
    """
    import math
    from lxml import etree

    geom = record.get("geom", {})
    x = px_to_emu(geom.get("x", 0) * scale_x)
    y = px_to_emu(geom.get("y", 0) * scale_y)
    w = px_to_emu(geom.get("w", 100) * scale_x)
    h = px_to_emu(geom.get("h", 80) * scale_y)

    shape = slide.shapes.add_shape(MSO_SHAPE.CUBE, x, y, w, h)

    # ── Depth adjustment ──
    adjust1 = geom.get("adjust1", 30)   # depth in pixels
    adjust2 = geom.get("adjust2", 135)  # angle in degrees

    # Compute effective depth (same clamp logic as canvas _effective_depth)
    raw_w = geom.get("w", 100)
    raw_h = geom.get("h", 80)
    depth = adjust1
    rad = math.radians(adjust2)
    abs_sin = abs(math.sin(rad))
    abs_cos = abs(math.cos(rad))
    if abs_sin > 1e-9:
        depth = min(depth, raw_w / abs_sin)
    if abs_cos > 1e-9:
        depth = min(depth, raw_h / abs_cos)
    depth = max(0.0, depth)

    max_dim = max(raw_w, raw_h)
    depth_ratio = min(0.5, depth / max_dim) if max_dim > 0 else 0.25
    try:
        shape.adjustments[0] = depth_ratio
    except (IndexError, AttributeError):
        pass

    # ── Map angle quadrant to flipH / flipV ──
    # Our angle: 0°=up, CW. PPTX CUBE default extrusion is upper-right
    # (front face lower-left), which corresponds to our 180–270° range.
    angle = adjust2 % 360
    xfrm = shape._element.spPr.xfrm
    if angle < 90:         # NE outward → flipH + flipV
        xfrm.set("flipH", "1")
        xfrm.set("flipV", "1")
    elif angle < 180:      # SE outward → flipH only
        xfrm.set("flipH", "1")
    elif angle < 270:      # SW outward → default (no flip)
        pass
    else:                  # NW outward → flipV only
        xfrm.set("flipV", "1")

    # ── Styles and text ──
    style = record.get("style", {})
    _apply_line_style(shape.line, style)
    _apply_fill_style(shape.fill, style)

    _add_text_to_shape(shape, record, style.get("text", {}))
    _apply_rotation(shape, geom)


def _add_orthocurve(slide, record: Dict[str, Any], scale_x: float, scale_y: float):
    """Add an orthogonal curve (M/H/V polyline) to the slide.

    Resolves relative M/H/V/Z nodes to absolute pixel coordinates, then
    builds a freeform polyline.  Corner bend radius (adjust1) is not
    exported — PowerPoint freeforms don't support native corner rounding.

    Args:
        slide: The PowerPoint slide
        record: Annotation record dict
        scale_x: Horizontal scale factor
        scale_y: Vertical scale factor
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    geom = record.get("geom", {})
    x = geom.get("x", 0) * scale_x
    y = geom.get("y", 0) * scale_y
    w = geom.get("w", 100) * scale_x
    h = geom.get("h", 100) * scale_y

    nodes = geom.get("nodes", [])
    if not nodes or len(nodes) < 2:
        return

    # ── Resolve M/H/V to absolute relative-coordinates ──
    points: List[Tuple[float, float]] = []
    prev_x, prev_y = 0.0, 0.0
    for node in nodes:
        cmd = node.get("cmd", "M")
        if cmd == "Z":
            continue
        if cmd == "M":
            prev_x, prev_y = node.get("x", 0.0), node.get("y", 0.0)
        elif cmd == "H":
            prev_x = node.get("x", prev_x)
        elif cmd == "V":
            prev_y = node.get("y", prev_y)
        points.append((prev_x, prev_y))

    if len(points) < 2:
        return

    # Convert to absolute EMU positions
    emu_points = [
        (px_to_emu(x + px * w), px_to_emu(y + py * h))
        for px, py in points
    ]

    # Build freeform polyline
    start_x, start_y = emu_points[0]
    builder = slide.shapes.build_freeform(start_x, start_y)
    builder.add_line_segments(emu_points[1:], close=False)
    shape = builder.convert_to_shape()

    _apply_rotation(shape, geom)

    # ── Styles ──
    style = record.get("style", {})
    _apply_line_style(shape.line, style)
    # No fill for open polylines
    shape.fill.background()

    # ── Arrowheads (same XML approach as _add_curve) ──
    arrow_mode = style.get("arrow", "none")
    if arrow_mode != "none":
        sp = shape._element
        spPr = sp.find(qn("a:spPr")) if sp.find(qn("a:spPr")) is not None else sp.spPr
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = etree.SubElement(spPr, qn("a:ln"))
        # headEnd must come before tailEnd in OOXML element order
        if arrow_mode in ("start", "both"):
            head_end = ln.find(qn("a:headEnd"))
            if head_end is None:
                head_end = etree.SubElement(ln, qn("a:headEnd"))
            head_end.set("type", "triangle")
        if arrow_mode in ("end", "both"):
            tail_end = ln.find(qn("a:tailEnd"))
            if tail_end is None:
                tail_end = etree.SubElement(ln, qn("a:tailEnd"))
            tail_end.set("type", "triangle")

    # ── Text label at polyline midpoint ──
    mid_idx = len(points) // 2
    prev_idx = max(0, mid_idx - 1)
    mid_rx = (points[prev_idx][0] + points[mid_idx][0]) / 2
    mid_ry = (points[prev_idx][1] + points[mid_idx][1]) / 2
    mid_x = geom.get("x", 0) + mid_rx * geom.get("w", 0)
    mid_y = geom.get("y", 0) + mid_ry * geom.get("h", 0)
    _add_line_curve_textbox(slide, record, geom, scale_x, scale_y,
                            mid_x=mid_x, mid_y=mid_y)


def export_to_pptx(
    annotations: List[Dict[str, Any]],
    output_path: str,
    scene_rect: Optional[Dict[str, float]] = None,
    background_png: Optional[str] = None,
) -> None:
    """
    Export annotations to a PowerPoint file.

    Args:
        annotations: List of annotation records from canvas items
        output_path: Path to save the .pptx file
        scene_rect: Scene rectangle with x, y, w, h (for scaling)
        background_png: Optional path to background image
    """
    prs = Presentation()

    # Set slide dimensions based on scene rect
    if scene_rect:
        scene_w = scene_rect.get("w", 800)
        scene_h = scene_rect.get("h", 600)
        prs.slide_width = px_to_emu(scene_w)
        prs.slide_height = px_to_emu(scene_h)
    else:
        scene_w = 800
        scene_h = 600

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Calculate scale factors (1:1 if using scene dimensions)
    scale_x = 1.0
    scale_y = 1.0

    # Add background image if provided
    if background_png:
        try:
            from pptx.util import Inches
            import os
            if os.path.exists(background_png):
                slide.shapes.add_picture(
                    background_png,
                    0, 0,
                    prs.slide_width,
                    prs.slide_height
                )
        except Exception:
            pass  # Skip if image can't be added

    # Flatten groups recursively before export
    def _flatten_annotations(anns):
        flat = []
        for rec in anns:
            if rec.get("kind") == "group":
                flat.extend(_flatten_annotations(rec.get("children", [])))
            else:
                flat.append(rec)
        return flat

    annotations = _flatten_annotations(annotations)

    # Sort annotations by z-index to maintain layering
    sorted_annotations = sorted(
        annotations,
        key=lambda a: a.get("z", 0)
    )

    # Add each shape to the slide
    for record in sorted_annotations:
        kind = record.get("kind", "")

        if kind == "rect":
            _add_rectangle(slide, record, scale_x, scale_y)
        elif kind == "roundedrect":
            _add_rounded_rectangle(slide, record, scale_x, scale_y)
        elif kind == "ellipse":
            _add_ellipse(slide, record, scale_x, scale_y)
        elif kind == "line":
            _add_line(slide, record, scale_x, scale_y)
        elif kind == "text":
            _add_text(slide, record, scale_x, scale_y)
        elif kind == "hexagon":
            _add_hexagon(slide, record, scale_x, scale_y)
        elif kind == "cylinder":
            _add_cylinder(slide, record, scale_x, scale_y)
        elif kind == "blockarrow":
            _add_blockarrow(slide, record, scale_x, scale_y)
        elif kind == "polygon":
            _add_polygon(slide, record, scale_x, scale_y)
        elif kind == "curve":
            _add_curve(slide, record, scale_x, scale_y)
        elif kind == "isocube":
            _add_isocube(slide, record, scale_x, scale_y)
        elif kind == "orthocurve":
            _add_orthocurve(slide, record, scale_x, scale_y)
        elif kind == "seqblock":
            _add_seqblock(slide, record, scale_x, scale_y)

    prs.save(output_path)
